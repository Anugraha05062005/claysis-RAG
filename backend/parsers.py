import os
import re
import csv
import json
import sqlite3
import base64
import logging
import email
from email import policy
import xml.etree.ElementTree as ET
from bs4 import BeautifulSoup
import pypdf
import docx
import pptx
import cv2
from PIL import Image
import requests

logger = logging.getLogger("parsers")

def encode_image_base64(image_path):
    """Encode image to base64 for API requests."""
    with open(image_path, "rb") as image_file:
        return base64.b64encode(image_file.read()).decode('utf-8')

class UnifiedParser:
    def __init__(self, groq_api_key=None):
        self.groq_api_key = groq_api_key

    def parse_file(self, file_path):
        """Auto-detect file extension and parse content into unified format."""
        ext = os.path.splitext(file_path)[1].lower()
        logger.info(f"Parsing file: {file_path} (Extension: {ext})")
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")

        if ext == '.pdf':
            return self.parse_pdf(file_path)
        elif ext == '.docx':
            return self.parse_docx(file_path)
        elif ext == '.pptx':
            return self.parse_pptx(file_path)
        elif ext in ('.txt', '.text'):
            return self.parse_txt(file_path)
        elif ext in ('.md', '.markdown'):
            return self.parse_markdown(file_path)
        elif ext in ('.html', '.htm'):
            return self.parse_html(file_path)
        elif ext == '.json':
            return self.parse_json(file_path)
        elif ext == '.xml':
            return self.parse_xml(file_path)
        elif ext == '.csv':
            return self.parse_csv(file_path)
        elif ext in ('.db', '.sqlite', '.sqlite3'):
            return self.parse_sqlite(file_path)
        elif ext in ('.eml', '.msg'):
            return self.parse_eml(file_path)
        elif ext in ('.png', '.jpg', '.jpeg', '.webp', '.bmp'):
            return self.parse_image(file_path)
        elif ext in ('.mp3', '.wav', '.m4a', '.ogg', '.flac'):
            return self.parse_audio(file_path)
        elif ext in ('.mp4', '.avi', '.mkv', '.mov', '.webm'):
            return self.parse_video(file_path)
        else:
            # Fallback to plain text read
            try:
                return self.parse_txt(file_path)
            except Exception as e:
                raise ValueError(f"Unsupported file format '{ext}' and failed to read as text: {e}")

    # 1. PDF Parser
    def parse_pdf(self, file_path):
        chunks = []
        filename = os.path.basename(file_path)
        
        with open(file_path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for idx, page in enumerate(reader.pages):
                text = page.extract_text()
                if text and text.strip():
                    chunks.append({
                        "text": text.strip(),
                        "type": "paragraph",
                        "metadata": {
                            "heading": f"Page {idx + 1}",
                            "page": idx + 1
                        }
                    })
        return {
            "source_type": "pdf",
            "source": filename,
            "chunks": chunks
        }

    # 2. DOCX Parser
    def parse_docx(self, file_path):
        chunks = []
        filename = os.path.basename(file_path)
        doc = docx.Document(file_path)
        
        current_heading = "Document Body"
        paragraphs = []
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Identify headings based on style name
            if para.style.name.startswith('Heading'):
                if paragraphs:
                    chunks.append({
                        "text": "\n".join(paragraphs),
                        "type": "paragraph",
                        "metadata": {"heading": current_heading}
                    })
                    paragraphs = []
                current_heading = text
            else:
                paragraphs.append(text)
                
        if paragraphs:
            chunks.append({
                "text": "\n".join(paragraphs),
                "type": "paragraph",
                "metadata": {"heading": current_heading}
            })
            
        # Parse tables in DOCX
        for idx, table in enumerate(doc.tables):
            table_data = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                table_data.append(" | ".join(cells))
            if table_data:
                chunks.append({
                    "text": "\n".join(table_data),
                    "type": "table",
                    "metadata": {"heading": f"Table {idx + 1}"}
                })
                
        return {
            "source_type": "docx",
            "source": filename,
            "chunks": chunks
        }

    # 3. PPTX Parser
    def parse_pptx(self, file_path):
        chunks = []
        filename = os.path.basename(file_path)
        prs = pptx.Presentation(file_path)
        
        for idx, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())
            
            if slide_texts:
                chunks.append({
                    "text": "\n".join(slide_texts),
                    "type": "paragraph",
                    "metadata": {
                        "heading": f"Slide {idx + 1}",
                        "slide_number": idx + 1
                    }
                })
        return {
            "source_type": "pptx",
            "source": filename,
            "chunks": chunks
        }

    # 4. TXT Parser
    def parse_txt(self, file_path):
        filename = os.path.basename(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read().strip()
            
        # We put the entire document in one chunk, and the chunker will split it.
        return {
            "source_type": "txt",
            "source": filename,
            "chunks": [{
                "text": text,
                "type": "paragraph",
                "metadata": {"heading": "Content"}
            }]
        }

    # 5. Markdown Parser
    def parse_markdown(self, file_path):
        filename = os.path.basename(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read().strip()
            
        return {
            "source_type": "markdown",
            "source": filename,
            "chunks": [{
                "text": text,
                "type": "paragraph",
                "metadata": {"heading": "Content"}
            }]
        }

    # 6. HTML Parser (Local Files)
    def parse_html(self, file_path):
        filename = os.path.basename(file_path)
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            html_content = f.read()
            
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Remove script and style elements
        for script_or_style in soup(["script", "style", "nav", "footer", "header"]):
            script_or_style.decompose()
            
        title = soup.title.string.strip() if soup.title else "HTML Document"
        
        # Simple extraction of text groups
        paragraphs = []
        for p in soup.find_all(['p', 'h1', 'h2', 'h3', 'h4', 'li']):
            txt = p.get_text().strip()
            if txt:
                paragraphs.append(txt)
                
        return {
            "source_type": "html",
            "source": filename,
            "chunks": [{
                "text": "\n".join(paragraphs),
                "type": "paragraph",
                "metadata": {
                    "heading": title
                }
            }]
        }

    # 7. JSON Parser (Extract Key-value Structures)
    def parse_json(self, file_path):
        filename = os.path.basename(file_path)
        with open(file_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            
        # Recursive flattener for JSON key-value pairs
        flat_pairs = []
        def flatten(obj, prefix=""):
            if isinstance(obj, dict):
                for k, v in obj.items():
                    flatten(v, f"{prefix}{k}.")
            elif isinstance(obj, list):
                for idx, item in enumerate(obj):
                    flatten(item, f"{prefix}[{idx}].")
            else:
                flat_pairs.append(f"{prefix[:-1]}: {obj}")
                
        flatten(data)
        
        return {
            "source_type": "json",
            "source": filename,
            "chunks": [{
                "text": "\n".join(flat_pairs),
                "type": "paragraph",
                "metadata": {"heading": "Key-Values"}
            }]
        }

    # 8. XML Parser
    def parse_xml(self, file_path):
        filename = os.path.basename(file_path)
        tree = ET.parse(file_path)
        root = tree.getroot()
        
        flat_elements = []
        def traverse(node, path=""):
            current_path = f"{path}/{node.tag}" if path else node.tag
            if node.text and node.text.strip():
                flat_elements.append(f"{current_path}: {node.text.strip()}")
            for attr, val in node.attrib.items():
                flat_elements.append(f"{current_path}@{attr}: {val}")
            for child in node:
                traverse(child, current_path)
                
        traverse(root)
        
        return {
            "source_type": "xml",
            "source": filename,
            "chunks": [{
                "text": "\n".join(flat_elements),
                "type": "paragraph",
                "metadata": {"heading": "Hierarchical Data"}
            }]
        }

    # 9. CSV Parser
    def parse_csv(self, file_path):
        filename = os.path.basename(file_path)
        rows = []
        with open(file_path, mode='r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for idx, row in enumerate(reader):
                if row:
                    rows.append(" , ".join(row))
                    
        return {
            "source_type": "csv",
            "source": filename,
            "chunks": [{
                "text": "\n".join(rows),
                "type": "table",
                "metadata": {"heading": "CSV Sheet"}
            }]
        }

    # 10. SQLite Parser
    def parse_sqlite(self, file_path):
        filename = os.path.basename(file_path)
        chunks = []
        
        conn = sqlite3.connect(file_path)
        cursor = conn.cursor()
        
        # Get list of tables
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
        tables = [t[0] for t in cursor.fetchall() if t[0] != 'sqlite_sequence']
        
        for table in tables:
            # Get table schema
            cursor.execute(f"PRAGMA table_info({table});")
            columns = [col[1] for col in cursor.fetchall()]
            
            # Fetch all rows
            cursor.execute(f"SELECT * FROM {table} LIMIT 100;")  # Safe limit
            rows = cursor.fetchall()
            
            table_string_lines = []
            table_string_lines.append(" | ".join(columns))
            table_string_lines.append("---" * len(columns))
            for row in rows:
                table_string_lines.append(" | ".join(str(val) for val in row))
                
            chunks.append({
                "text": "\n".join(table_string_lines),
                "type": "table",
                "metadata": {
                    "heading": f"Table: {table}",
                    "table_name": table
                }
            })
            
        conn.close()
        return {
            "source_type": "sqlite",
            "source": filename,
            "chunks": chunks
        }

    # 11. EML Parser (Emails)
    def parse_eml(self, file_path):
        filename = os.path.basename(file_path)
        with open(file_path, 'rb') as f:
            msg = email.message_from_binary_file(f, policy=policy.default)
            
        subject = msg.get('Subject', '(No Subject)')
        sender = msg.get('From', '(Unknown Sender)')
        recipient = msg.get('To', '(Unknown Recipient)')
        date = msg.get('Date', '(Unknown Date)')
        
        email_header = f"Subject: {subject}\nFrom: {sender}\nTo: {recipient}\nDate: {date}\n\n"
        
        body = ""
        if msg.is_multipart():
            for part in msg.walk():
                content_type = part.get_content_type()
                content_disposition = str(part.get_content_disposition())
                if content_type == 'text/plain' and 'attachment' not in content_disposition:
                    body += part.get_payload(decode=True).decode('utf-8', errors='ignore')
        else:
            body = msg.get_payload(decode=True).decode('utf-8', errors='ignore')
            
        full_content = email_header + body
        
        return {
            "source_type": "email",
            "source": filename,
            "chunks": [{
                "text": full_content,
                "type": "paragraph",
                "metadata": {
                    "heading": "Email content",
                    "subject": subject,
                    "from": sender,
                    "date": date
                }
            }]
        }

    # 12. Image OCR Parser
    def parse_image(self, file_path):
        filename = os.path.basename(file_path)
        extracted_text = ""
        
        if self.groq_api_key:
            try:
                base64_image = encode_image_base64(file_path)
                headers = {
                    "Authorization": f"Bearer {self.groq_api_key}",
                    "Content-Type": "application/json"
                }
                payload = {
                    "model": "llama-3.2-11b-vision-preview",
                    "messages": [
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "text",
                                    "text": "Perform OCR on this image. Extract and transcribe all text visible in the image exactly. Do not describe the image, just return the text."
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:image/jpeg;base64,{base64_image}"
                                    }
                                }
                            ]
                        }
                    ]
                }
                response = requests.post("https://api.groq.com/openapi/v1/chat/completions", headers=headers, json=payload, timeout=20)
                if response.status_code == 200:
                    data = response.json()
                    extracted_text = data["choices"][0]["message"]["content"].strip()
                else:
                    logger.warning(f"Groq Image OCR API failed with status {response.status_code}: {response.text}")
            except Exception as e:
                logger.error(f"Error calling Groq Vision API for OCR: {e}")
                
        # Local OCR Fallback (if Groq fails or is not configured)
        if not extracted_text:
            try:
                import pytesseract
                extracted_text = pytesseract.image_to_string(Image.open(file_path)).strip()
            except Exception as e:
                logger.warning(f"Pytesseract local OCR failed (probably not installed): {e}")
                
        if not extracted_text:
            # Fallback placeholder text if OCR failed
            extracted_text = f"[OCR Extracted Metadata for image '{filename}']\n(Please provide a Groq API Key or install Tesseract OCR for text extraction)."

        return {
            "source_type": "image",
            "source": filename,
            "chunks": [{
                "text": extracted_text,
                "type": "ocr",
                "metadata": {
                    "heading": "OCR Text Extraction"
                }
            }]
        }

    # 13. Audio Parser (Whisper)
    def parse_audio(self, file_path):
        filename = os.path.basename(file_path)
        transcript = ""
        
        if self.groq_api_key:
            try:
                headers = {
                    "Authorization": f"Bearer {self.groq_api_key}"
                }
                files = {
                    "file": (filename, open(file_path, "rb"), "audio/mpeg"),
                    "model": (None, "whisper-large-v3"),
                    "response_format": (None, "verbose_json")
                }
                response = requests.post("https://api.groq.com/openapi/v1/audio/transcriptions", headers=headers, files=files, timeout=40)
                if response.status_code == 200:
                    data = response.json()
                    # Groq verbose_json returns a dict containing segments
                    segments = data.get("segments", [])
                    chunks = []
                    for seg in segments:
                        start = seg.get("start", 0)
                        minutes = int(start // 60)
                        seconds = int(start % 60)
                        timestamp_str = f"{minutes:02d}:{seconds:02d}"
                        
                        chunks.append({
                            "text": seg.get("text", "").strip(),
                            "type": "transcript",
                            "metadata": {
                                "heading": f"Audio [{timestamp_str}]",
                                "timestamp": timestamp_str
                            }
                        })
                    return {
                        "source_type": "audio",
                        "source": filename,
                        "chunks": chunks
                    }
                else:
                    logger.warning(f"Groq Whisper API failed with status {response.status_code}: {response.text}")
            except Exception as e:
                logger.error(f"Error calling Groq Whisper API for audio: {e}")
                
        # Local whisper fallback
        try:
            import whisper
            model = whisper.load_model("tiny")
            result = model.transcribe(file_path)
            chunks = []
            for seg in result.get("segments", []):
                start = seg.get("start", 0)
                minutes = int(start // 60)
                seconds = int(start % 60)
                timestamp_str = f"{minutes:02d}:{seconds:02d}"
                
                chunks.append({
                    "text": seg.get("text", "").strip(),
                    "type": "transcript",
                    "metadata": {
                        "heading": f"Audio [{timestamp_str}]",
                        "timestamp": timestamp_str
                    }
                })
            return {
                "source_type": "audio",
                "source": filename,
                "chunks": chunks
            }
        except Exception as e:
            logger.warning(f"Local Whisper transcribing failed: {e}")
            
        # Standard fallback placeholder
        return {
            "source_type": "audio",
            "source": filename,
            "chunks": [{
                "text": f"[Audio file uploaded: '{filename}']. Audio transcription is not available locally. Please configure a Groq API Key for instant Whisper cloud transcription.",
                "type": "transcript",
                "metadata": {
                    "heading": "Audio Transcript Placeholder",
                    "timestamp": "00:00"
                }
            }]
        }

    # 14. Video Parser (Timeline OCR + Audio Transcribe)
    def parse_video(self, file_path):
        filename = os.path.basename(file_path)
        chunks = []
        
        # 1. Timeline Frame OCR Extraction using OpenCV
        try:
            cap = cv2.VideoCapture(file_path)
            fps = cap.get(cv2.CAP_PROP_FPS)
            total_frames = cap.get(cv2.CAP_PROP_FRAME_COUNT)
            duration_sec = total_frames / fps if fps > 0 else 0
            
            # Sample frames every 10 seconds
            interval_sec = 10
            frame_step = int(fps * interval_sec) if fps > 0 else 100
            
            current_frame = 0
            while current_frame < total_frames:
                cap.set(cv2.CAP_PROP_POS_FRAMES, current_frame)
                ret, frame = cap.read()
                if not ret:
                    break
                    
                timestamp_sec = current_frame / fps if fps > 0 else 0
                minutes = int(timestamp_sec // 60)
                seconds = int(timestamp_sec % 60)
                timestamp_str = f"{minutes:02d}:{seconds:02d}"
                
                # Save frame temporarily
                temp_frame_path = f"temp_frame_{timestamp_str.replace(':', '_')}.jpg"
                cv2.imwrite(temp_frame_path, frame)
                
                # Perform OCR on frame
                ocr_text = ""
                if self.groq_api_key:
                    try:
                        base64_image = encode_image_base64(temp_frame_path)
                        headers = {
                            "Authorization": f"Bearer {self.groq_api_key}",
                            "Content-Type": "application/json"
                        }
                        payload = {
                            "model": "llama-3.2-11b-vision-preview",
                            "messages": [
                                {
                                    "role": "user",
                                    "content": [
                                        {
                                            "type": "text",
                                            "text": "Extract and transcribe all text visible in this video slide frame exactly. If no text is present, return nothing."
                                        },
                                        {
                                            "type": "image_url",
                                            "image_url": {
                                                "url": f"data:image/jpeg;base64,{base64_image}"
                                            }
                                        }
                                    ]
                                }
                            ]
                        }
                        response = requests.post("https://api.groq.com/openapi/v1/chat/completions", headers=headers, json=payload, timeout=20)
                        if response.status_code == 200:
                            data = response.json()
                            ocr_text = data["choices"][0]["message"]["content"].strip()
                    except Exception:
                        pass
                        
                if not ocr_text:
                    # Simple fallback
                    ocr_text = f"[Visual Frame Content at {timestamp_str}]"
                    
                # Clean temp file
                if os.path.exists(temp_frame_path):
                    os.remove(temp_frame_path)
                    
                if ocr_text:
                    chunks.append({
                        "text": ocr_text,
                        "type": "ocr",
                        "metadata": {
                            "heading": f"Video Frame [{timestamp_str}]",
                            "timestamp": timestamp_str
                        }
                    })
                    
                current_frame += frame_step
            cap.release()
        except Exception as e:
            logger.warning(f"Failed to extract video frames via OpenCV: {e}")
            
        # 2. Audio Transcription (Whisper)
        # Attempt to upload video file directly to Whisper API (supported by Groq up to 25MB)
        audio_chunks = []
        audio_success = False
        if self.groq_api_key and os.path.getsize(file_path) < 25 * 1024 * 1024:
            try:
                headers = {"Authorization": f"Bearer {self.groq_api_key}"}
                files = {
                    "file": (filename, open(file_path, "rb"), "audio/mpeg"),
                    "model": (None, "whisper-large-v3"),
                    "response_format": (None, "verbose_json")
                }
                response = requests.post("https://api.groq.com/openapi/v1/audio/transcriptions", headers=headers, files=files, timeout=60)
                if response.status_code == 200:
                    data = response.json()
                    segments = data.get("segments", [])
                    for seg in segments:
                        start = seg.get("start", 0)
                        minutes = int(start // 60)
                        seconds = int(start % 60)
                        timestamp_str = f"{minutes:02d}:{seconds:02d}"
                        
                        audio_chunks.append({
                            "text": seg.get("text", "").strip(),
                            "type": "transcript",
                            "metadata": {
                                "heading": f"Audio Speech [{timestamp_str}]",
                                "timestamp": timestamp_str
                            }
                        })
                    audio_success = True
            except Exception as e:
                logger.error(f"Error calling Groq Whisper for video file: {e}")
                
        # If Groq failed or file is too large, try local whisper (if available)
        if not audio_success:
            try:
                # Local whisper handles video files directly (if ffmpeg libraries are present in path)
                import whisper
                model = whisper.load_model("tiny")
                result = model.transcribe(file_path)
                for seg in result.get("segments", []):
                    start = seg.get("start", 0)
                    minutes = int(start // 60)
                    seconds = int(start % 60)
                    timestamp_str = f"{minutes:02d}:{seconds:02d}"
                    
                    audio_chunks.append({
                        "text": seg.get("text", "").strip(),
                        "type": "transcript",
                        "metadata": {
                            "heading": f"Audio Speech [{timestamp_str}]",
                            "timestamp": timestamp_str
                        }
                    })
                audio_success = True
            except Exception as e:
                logger.warning(f"Local Whisper transcribing for video failed: {e}")
                
        if not audio_success:
            audio_chunks.append({
                "text": f"[Video Audio Speech Transcript for '{filename}']\n(Cloud Groq API Key required or local whisper/ffmpeg must be installed to extract voice script).",
                "type": "transcript",
                "metadata": {
                    "heading": "Audio Transcript Placeholder",
                    "timestamp": "00:00"
                }
            })
            
        # Combine OCR chunks + Audio speech chunks
        chunks.extend(audio_chunks)
        
        # Sort chunks by timestamp if possible
        def get_sec(chunk):
            ts = chunk.get("metadata", {}).get("timestamp", "00:00")
            try:
                parts = ts.split(":")
                return int(parts[0]) * 60 + int(parts[1])
            except Exception:
                return 0
                
        chunks = sorted(chunks, key=get_sec)
        
        return {
            "source_type": "video",
            "source": filename,
            "chunks": chunks
        }
